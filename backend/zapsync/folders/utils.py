import jwt
import os
import re
import io
import json
import csv
import fitz  # PyMuPDF
import zipfile
import pytesseract
import logging
from PIL import Image
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from better_profanity import profanity
from django.conf import settings
from rest_framework.exceptions import AuthenticationFailed
from django.core.files.uploadedfile import UploadedFile as DjangoUploadedFile
from django.core.exceptions import ValidationError
import re
# Set up logging
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.DEBUG)

# Load profanity words
profanity.load_censor_words()
# logger.debug(f"[Profanity Init] Profanity word list loaded with {len(profanity.CENSOR_WORDSET)} words")

def decode_jwt(token):
    try:
        decoded = jwt.decode(token, settings.SECRET_KEY, algorithms=["HS256"])
        # logger.debug(f"[JWT] Decoded user_id: {decoded.get('user_id')}")
        return decoded['user_id']
    except jwt.ExpiredSignatureError:
        logger.error("[JWT] Token expired")
        raise AuthenticationFailed('Token expired')
    except jwt.InvalidTokenError:
        logger.error("[JWT] Invalid token")
        raise AuthenticationFailed('Invalid token')


def extract_text_from_file(file: DjangoUploadedFile, file_name: str) -> str:
    ext = os.path.splitext(file_name)[1].lower()
    text = ""
    # logger.debug(f"[Extract] Attempting to extract from {file_name} with extension {ext}")

    try:
        if ext == ".txt":
            text = file.read().decode("utf-8", errors="ignore")

        elif ext == ".pdf":
            with fitz.open(stream=file.read(), filetype="pdf") as doc:
                text = " ".join(page.get_text() for page in doc)

        elif ext == ".docx":
            doc = Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])

        elif ext == ".pptx":
            prs = Presentation(file)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

        elif ext == ".xlsx":
            wb = load_workbook(file)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"

        elif ext == ".csv":
            text = file.read().decode("utf-8", errors="ignore")

        elif ext == ".json":
            json_data = json.loads(file.read().decode("utf-8", errors="ignore"))
            text = json.dumps(json_data)

        elif ext == ".zip":
            zip_text = ""
            with zipfile.ZipFile(file) as archive:
                for name in archive.namelist():
                    with archive.open(name) as zipped_file:
                        zip_text += extract_text_from_file(zipped_file, name)
            text = zip_text

        # logger.debug(f"[Extract] Extracted {len(text)} characters of text from {file_name}")

    except Exception as e:
        logger.error(f"[Extract] Failed to extract text from {file_name}: {e}")

    return text


def sanitize_text(text: str) -> str:
    # logger.debug(f"[Sanitize] Raw length: {len(text)}")
    text = re.sub(r'[^a-zA-Z\s]', ' ', text)
    cleaned = re.sub(r'\s+', ' ', text).lower().strip()
    # logger.debug(f"[Sanitize] Cleaned length: {len(cleaned)}")
    return cleaned


def check_for_profanity(file: DjangoUploadedFile, file_name: str) -> bool:
    try:
        # logger.debug(f"[Profanity] Checking file: {file_name}")
        text = extract_text_from_file(file, file_name)
        file.seek(0)
        clean_text = sanitize_text(text)

        # Get profane words (not just True/False)
        profane_words = [
            word for word in clean_text.split() if profanity.contains_profanity(word)
        ]

        # logger.debug(f"[Profanity] Detected words: {profane_words}")

        return bool(profane_words)

    except Exception as e:
        logger.error(f"[Profanity] Error while checking for profanity: {e}")
        return False




def validate_folder_name(value):
    """
    Validates that folder names:
    - Are not empty
    - Don't contain special characters
    - Are not too long
    """
    if not value.strip():
        raise ValidationError("Folder name cannot be empty")
    
    if len(value) > 255:
        raise ValidationError("Folder name is too long (max 255 characters)")
    